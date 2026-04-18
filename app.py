# ================================================================
# DIABETES AI AGENT — Full Platform
# 6 Modules: Insulin Delivery, CGM, Diagnostics, Drugs,
#            Nutrition, Population Health
#
# Student: Hugo Silva | ID: 74964557
# ITU | SPRING 2026 | HCM 535 – Data Analytics in Healthcare
#
# HOW TO RUN LOCALLY:
#   pip install streamlit anthropic python-pptx
#   streamlit run app.py
#
# HOW TO DEPLOY:
#   Push to GitHub → share.streamlit.io
#   Add ANTHROPIC_API_KEY to Streamlit Secrets
# ================================================================

import streamlit as st
import anthropic
import json
import re
import os
from datetime import datetime
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── PAGE CONFIG ─────────────────────────────────────────────────
st.set_page_config(
    page_title="Diabetes AI Agent | Hugo Silva",
    page_icon="🩺",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ── CUSTOM CSS ───────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Syne:wght@400;600;700;800&family=DM+Serif+Display:ital@0;1&display=swap');

html, body, [class*="css"] { font-family: 'Syne', sans-serif; }

.main-header {
    background: linear-gradient(135deg, #042F2E 0%, #0F766E 100%);
    padding: 2rem 2.5rem;
    border-radius: 14px;
    margin-bottom: 1.5rem;
    border-left: 6px solid #F59E0B;
}
.main-header h1 { font-family: 'DM Serif Display', serif; font-size: 2.2rem; color: #CCFBF1; margin: 0 0 0.3rem 0; }
.main-header p { color: #5EEAD4; font-size: 0.9rem; margin: 0; opacity: 0.85; }
.student-badge {
    background: rgba(245,158,11,0.15);
    border: 1px solid #F59E0B;
    border-radius: 8px;
    padding: 0.5rem 1rem;
    color: #FCD34D;
    font-size: 0.8rem;
    margin-top: 0.8rem;
    display: inline-block;
}

.module-card {
    background: var(--background-color);
    border: 1px solid #E2E8F0;
    border-radius: 12px;
    padding: 1rem;
    margin-bottom: 0.6rem;
    border-left: 4px solid #14B8A6;
}
.module-card h4 { margin: 0 0 4px 0; font-size: 0.9rem; color: #0F766E; }
.module-card p { margin: 0; font-size: 0.78rem; color: #64748B; }

.metric-card {
    background: #F8FAFC;
    border: 1px solid #E2E8F0;
    border-radius: 12px;
    padding: 1.2rem;
    border-top: 3px solid #14B8A6;
    margin-bottom: 0.8rem;
}
.metric-card h3 { font-size: 2rem; font-weight: 800; color: #0F766E; margin: 0; font-family: 'DM Serif Display', serif; }
.metric-card p { font-size: 0.75rem; color: #64748B; margin: 0.2rem 0 0 0; }

.report-section { margin-bottom: 1.5rem; }
.section-header {
    background: #042F2E;
    color: #CCFBF1;
    padding: 0.5rem 1rem;
    border-radius: 8px 8px 0 0;
    font-size: 0.85rem;
    font-weight: 700;
    border-left: 4px solid #F59E0B;
}

.stDownloadButton button {
    background: #0F766E !important;
    color: white !important;
    border: none !important;
    border-radius: 8px !important;
    font-family: 'Syne', sans-serif !important;
    font-weight: 700 !important;
    padding: 0.6rem 1.5rem !important;
    font-size: 0.9rem !important;
    width: 100% !important;
}
.stDownloadButton button:hover { background: #14B8A6 !important; }

.stButton button {
    background: #042F2E !important;
    color: #CCFBF1 !important;
    border: 1px solid #0F766E !important;
    border-radius: 8px !important;
    font-family: 'Syne', sans-serif !important;
    font-weight: 700 !important;
    width: 100% !important;
    font-size: 0.95rem !important;
}
.stButton button:hover { background: #0F766E !important; color: white !important; }

.teal-divider { border: none; border-top: 2px solid #14B8A6; margin: 1rem 0; opacity: 0.3; }

.footer {
    text-align: center;
    padding: 1.5rem;
    color: #94A3B8;
    font-size: 0.75rem;
    border-top: 1px solid #E2E8F0;
    margin-top: 2rem;
}

.badge-success { background: #EAF3DE; color: #27500A; border-radius: 20px; padding: 3px 12px; font-size: 0.75rem; font-weight: 600; }
.badge-running { background: #FAEEDA; color: #633806; border-radius: 20px; padding: 3px 12px; font-size: 0.75rem; font-weight: 600; }
</style>
""", unsafe_allow_html=True)

# ── STUDENT INFO ─────────────────────────────────────────────────
STUDENT = {
    "name":   "Hugo Silva",
    "id":     "74964557",
    "school": "ITU  |  SPRING 2026",
    "course": "HCM 535 – Data Analytics Application in Healthcare",
}
FOOTER_TXT = f"{STUDENT['name']}  ·  ID {STUDENT['id']}  ·  {STUDENT['school']}  ·  {STUDENT['course']}"

# ── MODULE DEFINITIONS ────────────────────────────────────────────
MODULES = {
    "🩺 Full Diabetes AI Platform": {
        "icon": "🩺",
        "color": "#0F766E",
        "desc": "Complete overview of all AI applications in diabetes care",
        "tag": "Comprehensive"
    },
    "💉 AI Insulin Delivery (Smart Pumps)": {
        "icon": "💉",
        "color": "#0891B2",
        "desc": "Closed-loop AID systems, CGM integration, vendor comparison",
        "tag": "Clinical"
    },
    "🩸 AI Glucose Monitoring (CGM)": {
        "icon": "🩸",
        "color": "#DC2626",
        "desc": "Continuous glucose monitors, accuracy, non-invasive sensing",
        "tag": "Devices"
    },
    "🧠 AI Diagnostics & Early Detection": {
        "icon": "🧠",
        "color": "#7C3AED",
        "desc": "T2D prediction, retinopathy screening, complication risk AI",
        "tag": "Diagnostics"
    },
    "💊 AI Drug Management (GLP-1 & Meds)": {
        "icon": "💊",
        "color": "#B45309",
        "desc": "GLP-1 drugs, smart pens, adherence AI, dosing algorithms",
        "tag": "Pharmacology"
    },
    "🥗 AI Nutrition & Lifestyle Coaching": {
        "icon": "🥗",
        "color": "#16A34A",
        "desc": "Meal planning AI, carb counting, activity coaching, apps",
        "tag": "Lifestyle"
    },
    "📊 Population Health & Analytics": {
        "icon": "📊",
        "color": "#0F766E",
        "desc": "Predictive risk models, payer analytics, clinical trial AI",
        "tag": "Analytics"
    },
    "⚖️ Vendor Comparison Only": {
        "icon": "⚖️",
        "color": "#64748B",
        "desc": "Side-by-side comparison of all major diabetes AI vendors",
        "tag": "Comparison"
    },
    "🎯 Recommendation Only": {
        "icon": "🎯",
        "color": "#F59E0B",
        "desc": "Direct evidence-based recommendation by patient profile",
        "tag": "Decision"
    },
}

# ── PPTX PALETTE ─────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)
DARK    = rgb(0x04, 0x2F, 0x2E)
PRIMARY = rgb(0x0F, 0x76, 0x6E)
MID     = rgb(0x14, 0xB8, 0xA6)
LIGHT   = rgb(0xCC, 0xFB, 0xF1)
ACCENT  = rgb(0xF5, 0x9E, 0x0B)
WHITE   = rgb(0xFF, 0xFF, 0xFF)
OFFWHITE= rgb(0xF8, 0xFA, 0xFC)
GRAY    = rgb(0x64, 0x74, 0x8B)
LTGRAY  = rgb(0xE2, 0xE8, 0xF0)
TEXT    = rgb(0x0F, 0x17, 0x2A)
RED     = rgb(0xDC, 0x26, 0x26)
GREEN   = rgb(0x16, 0xA3, 0x4A)
PURPLE  = rgb(0x7C, 0x3A, 0xED)
BLUE    = rgb(0x08, 0x91, 0xB2)
FT = "Georgia"
FB = "Calibri"

# ── PPTX HELPERS ─────────────────────────────────────────────────
def add_rect(s, x, y, w, h, fill, line=None):
    shape = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line: shape.line.color.rgb = line; shape.line.width = Pt(0.5)
    else: shape.line.fill.background()
    return shape

def add_text(s, text, x, y, w, h, size=12, bold=False,
             italic=False, color=None, align=PP_ALIGN.LEFT, font_name="Calibri"):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = str(text)
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.name = font_name
    run.font.color.rgb = color if color else TEXT
    return tb

def hdr(s, title, sub=None):
    add_rect(s, 0, 0, 10, 0.62, DARK)
    add_rect(s, 0, 0.62, 10, 0.05, MID)
    add_text(s, title, 0.4, 0.07, 7, 0.52, 18, True, color=WHITE, font_name=FT)
    if sub:
        add_text(s, sub, 6.8, 0.07, 2.9, 0.52, 9, False, True, color=MID, align=PP_ALIGN.RIGHT)

def ftr(s, pg=None):
    add_rect(s, 0, 5.36, 10, 0.265, DARK)
    add_text(s, FOOTER_TXT, 0.3, 5.36, 8.8, 0.265, 7, color=MID, align=PP_ALIGN.CENTER)
    if pg:
        add_text(s, str(pg), 9.55, 5.36, 0.35, 0.265, 7, color=LIGHT, align=PP_ALIGN.RIGHT)

# ── SYSTEM PROMPTS PER MODULE ────────────────────────────────────
MODULE_SYSTEMS = {
    "🩺 Full Diabetes AI Platform": """You are an expert AI research agent covering ALL aspects of
AI in diabetes management. You have comprehensive knowledge of:

INSULIN DELIVERY: Medtronic MiniMed 780G (TIR 78%, 700K users), Insulet Omnipod 5
(TIR 69%, n=69,902, hypo 1.12%, T2D cleared 2024), Tandem Control-IQ+ (TIR 73.6%,
94% auto mode), Beta Bionics iLet (fully autonomous), CamAPS FX (FDA 2024).

CGM DEVICES: Dexcom G7 (MARD 8.1%, 10-day wear), Abbott FreeStyle Libre 3
(MARD 7.9%, 14-day, no fingerstick), Medtronic Guardian 4, Senseonics Eversense
(180-day implantable), non-invasive sensing research.

DIAGNOSTICS AI: Google's ARDA (diabetic retinopathy, 90%+ sensitivity), IDx-DR
(FDA cleared autonomous retinopathy screening), AI for neuropathy detection,
T2D prediction models (ADA risk score + ML achieving AUC 0.85+).

GLP-1 & DRUGS: Semaglutide (Ozempic/Wegovy), Tirzepatide (Mounjaro/Zepbound),
smart insulin pens (NovoPen 6, InPen by Medtronic), AI adherence platforms,
once-weekly insulin icodec.

NUTRITION AI: One Drop, Nutrino, DayTwo (microbiome-based), Glucose Goddess
approach, CGM-guided nutrition, meal planning AI apps.

POPULATION HEALTH: Komodo Health, Clarify Health, IBM Watson Health diabetes
modules, CMS diabetes prevention program AI, risk stratification models.

Always be specific with data. When returning JSON return ONLY raw JSON.""",

    "💉 AI Insulin Delivery (Smart Pumps)": """You are an expert in automated insulin delivery (AID) systems.
Clinical data 2024-2025:
- Medtronic MiniMed 780G: TIR 76-78%, GMI 6.8%, 700K+ users, AHCL algorithm,
  Guardian 4 CGM, SmartGuard
- Insulet Omnipod 5: TIR 69% (n=69,902 T1D), hypo 1.12% (lowest), T2D FDA 2024,
  tubeless, HbA1c -0.8% T2D trial
- Tandem Control-IQ+: TIR 73.6% (from 63.6%), 94% auto mode, 365K users, MPC
- Beta Bionics iLet: fully autonomous, no carb counting, strong pediatric data
- CamAPS FX: FDA cleared 2024, app-based, best pregnancy outcomes
When returning JSON return ONLY raw JSON.""",

    "🩸 AI Glucose Monitoring (CGM)": """You are an expert in continuous glucose monitoring (CGM) technology.
Key data 2024-2025:
- Dexcom G7: MARD 8.1%, 10-day wear, no warmup, direct to Apple Watch,
  FDA cleared, 60-min grace period, works with Omnipod 5 and Control-IQ
- Abbott FreeStyle Libre 3: MARD 7.9% (best in class), 14-day wear,
  no fingerstick calibration, real-time streaming, $75/month, widest coverage
- Medtronic Guardian 4: factory calibrated, integrates with 780G,
  predictive alerts 60 min ahead
- Senseonics Eversense E3: 180-day implantable, MARD 8.5%, vibration alerts,
  requires physician insertion
- Non-invasive research: Apple Watch glucose (2026 target), Samsung Galaxy Ring,
  GlucoLight optical sensing — none FDA cleared yet
- CGM in T2D: expanding coverage, Medicare approved 2023, clinical benefit shown
When returning JSON return ONLY raw JSON.""",

    "🧠 AI Diagnostics & Early Detection": """You are an expert in AI diagnostics for diabetes and its complications.
Key data 2024-2025:
- Diabetic Retinopathy: IDx-DR (FDA cleared, 87.2% sensitivity, 90.7% specificity,
  autonomous — no ophthalmologist needed), Google ARDA (90%+ sensitivity),
  EyeArt (Eyenuk, FDA cleared)
- T2D Prediction: Finnish Diabetes Risk Score + ML achieves AUC 0.89,
  NHS Diabetes Prevention Programme AI identified 2M at-risk patients,
  UK Biobank models predict T2D 10 years ahead
- Diabetic Neuropathy: AI gait analysis, thermal imaging AI, corneal confocal
  microscopy AI — reducing missed diagnoses by 40%
- Kidney Disease (DKD): AI urinalysis, eGFR prediction models, early CKD
  detection 3-5 years earlier than standard care
- Wound Care: Tissue Analytics, Swift Medical AI wound assessment reducing
  amputation risk by 30% with early detection
- General: 70% of clinical decisions rely on diagnostic tests; AI reducing
  diagnostic delays by up to 50% (JMIR 2026 systematic review)
When returning JSON return ONLY raw JSON.""",

    "💊 AI Drug Management (GLP-1 & Meds)": """You are an expert in AI-powered diabetes drug management.
Key data 2024-2025:
- GLP-1 Revolution: Semaglutide (Ozempic) — HbA1c reduction 1.5-1.8%,
  weight loss 15%; Tirzepatide (Mounjaro) — HbA1c reduction 2.0-2.3%,
  weight loss 20-22%, superior to all comparators in SURPASS trials
- Smart Insulin Pens: NovoPen 6 (Novo Nordisk) records doses, syncs to app;
  Lilly connected pen; InPen (Medtronic) — dose guidance, reminders, CGM sync
- Once-Weekly Insulin: Icodec (Awiqli) — approved EU/Canada, FDA CRL 2024,
  resubmission expected 2025; reduces injection burden by 86%
- AI Adherence Platforms: Voluntis Insulia (AI insulin titration, FDA cleared),
  Amalgam Rx DosePal, Livongo (acquired by Teladoc) — AI coaching
- Precision Dosing: pharmacogenomics AI matching drug to genetic profile,
  reducing adverse events by 25%
- Digital Therapeutics: Omada Health, Virta Health (T2D reversal program —
  60% of patients off insulin at 1 year)
When returning JSON return ONLY raw JSON.""",

    "🥗 AI Nutrition & Lifestyle Coaching": """You are an expert in AI-powered nutrition and lifestyle management for diabetes.
Key data 2024-2025:
- CGM-Guided Nutrition: DayTwo (microbiome + CGM, personalized glycemic response
  prediction, 80% reduction in glucose spikes in trials), Levels Health
  (continuous glucose + AI insights, $199/month)
- Meal Planning AI: One Drop (AI coaching + CGM, HbA1c reduction 0.6% at 6mo),
  Nutrisense (CGM + registered dietitian AI), Signos (weight loss + CGM AI)
- Carb Counting: Carb Manager AI (95% accuracy food recognition),
  MyFitnessPal diabetes mode, Noom diabetes program
- Physical Activity AI: Garmin Connect IQ diabetes features, Apple Health
  activity impact on glucose, Whoop recovery + glucose correlation
- Behavioral AI: Virta Health (T2D reversal — low carb + coaching AI,
  60% insulin-free at 1 year), Omada (CDC-recognized DPP, 4-7% weight loss),
  Noom Med (GLP-1 + behavioral AI)
- Key finding: Lifestyle AI interventions reduce T2D progression by 58%
  (comparable to metformin) in high-risk individuals
When returning JSON return ONLY raw JSON.""",

    "📊 Population Health & Analytics": """You are an expert in AI-powered population health management for diabetes.
Key data 2024-2025:
- Scale: 537 million people with diabetes globally (IDF 2021), 783M projected 2045;
  $966B annual healthcare spend; 50% undiagnosed globally
- Risk Stratification AI: Komodo Health (claims + clinical AI, identifies at-risk
  patients 18 months earlier), Clarify Health (ML risk scores in EHR workflow),
  Arcadia Analytics (population health dashboard used by 50+ health systems)
- CMS Programs: Medicare Diabetes Prevention Program (mDPP) — AI-enhanced
  delivery, 5-7% weight loss, $2,650 savings/patient/year
- Payer Analytics: UnitedHealth OptumIQ diabetes module, Aetna Health AI
  (reduced T2D hospitalizations 20%), BCBS diabetes management AI
- Clinical Trial AI: Medidata Rave AI for diabetes trials, TriNetX (real-world
  evidence platform), IBM Watson for oncology/diabetes crossover research
- Health Equity: AI bias in diabetes risk models — Black patients under-predicted
  by 30% using standard HbA1c; UCSF/Stanford equity-adjusted models in development
- ROI: Every $1 invested in AI diabetes prevention returns $3.50 in avoided costs
When returning JSON return ONLY raw JSON.""",

    "⚖️ Vendor Comparison Only": """You are an expert comparing all major diabetes AI vendors.
Provide detailed head-to-head comparison across all diabetes AI categories:
insulin pumps, CGM devices, diagnostic AI, drug management, nutrition AI.
Include: product name, FDA status, clinical outcomes, pricing, target patient,
key strength, key limitation.
When returning JSON return ONLY raw JSON.""",

    "🎯 Recommendation Only": """You are an expert providing direct diabetes AI recommendations.
Give specific, actionable recommendations by patient type.
Be direct — name the product, explain why, include clinical evidence.
Always tailor to the patient profile specified.
When returning JSON return ONLY raw JSON.""",
}

# ── REPORT PROMPTS ────────────────────────────────────────────────
def build_prompt(module, profile, depth):
    profile_map = {
        "General": "general diabetes patients",
        "Type 1 Diabetes": "Type 1 diabetes patients",
        "Type 2 Diabetes": "Type 2 diabetes patients",
        "Prediabetes / At-Risk": "prediabetes and at-risk individuals",
        "Pediatric": "pediatric diabetes patients (children and adolescents)",
        "Elderly (65+)": "elderly patients (65+) with diabetes",
        "Newly Diagnosed": "newly diagnosed diabetes patients",
        "Healthcare Provider": "endocrinologists and diabetes care teams",
    }
    depth_map = {
        "Executive Summary": "Be concise — 2-3 paragraphs per section. Focus on key findings.",
        "Detailed Analysis": "Be thorough. Include specific numbers, trial names, and clinical context.",
        "Clinical Deep-Dive": "Be highly technical. Include trial names, n-values, p-values, effect sizes, confidence intervals.",
    }
    p = profile_map.get(profile, "general diabetes patients")
    d = depth_map.get(depth, depth_map["Detailed Analysis"])

    module_prompts = {
        "🩺 Full Diabetes AI Platform": f"""Generate a comprehensive intelligence report on ALL AI applications
in diabetes management for {p}. {d}

Structure with these exact sections:
1. Executive Summary — state of AI in diabetes 2025, 3 key findings
2. AI Insulin Delivery — smart pumps, closed-loop systems, top vendors, clinical outcomes
3. AI Glucose Monitoring — CGM technology, top devices, accuracy data, non-invasive future
4. AI Diagnostics — retinopathy screening, T2D prediction, complication detection
5. AI Drug Management — GLP-1 revolution, smart pens, adherence AI
6. AI Nutrition & Lifestyle — CGM-guided nutrition, coaching apps, behavioral AI
7. Population Health AI — risk stratification, payer analytics, equity considerations
8. Overall Recommendation — best AI ecosystem for {p}
9. Future Outlook 2025-2030
10. Risks & Limitations across all categories""",

        "💉 AI Insulin Delivery (Smart Pumps)": f"""Generate a comprehensive report on AI smart insulin pumps for {p}. {d}
Sections: Executive Summary, How AI Works (closed-loop, CGM, algorithms),
Clinical Effectiveness (TIR, HbA1c, hypoglycemia data), Vendor Comparison
(Medtronic 780G, Omnipod 5, Control-IQ+, iLet, CamAPS FX), Recommendation,
Risks & Limitations, Future Outlook.""",

        "🩸 AI Glucose Monitoring (CGM)": f"""Generate a comprehensive report on AI-powered glucose monitoring
for {p}. {d}
Sections: Executive Summary, How CGM + AI Works, Clinical Evidence
(accuracy data, MARD, outcomes), Device Comparison (Dexcom G7, FreeStyle Libre 3,
Guardian 4, Eversense E3), Non-Invasive Future, Recommendation, Costs & Access.""",

        "🧠 AI Diagnostics & Early Detection": f"""Generate a comprehensive report on AI diagnostics for diabetes
and its complications for {p}. {d}
Sections: Executive Summary, Retinopathy AI (IDx-DR, Google ARDA),
T2D Prediction Models, Neuropathy & Wound Care AI, Kidney Disease Detection,
Vendor Comparison, Clinical Evidence, Recommendation, Limitations.""",

        "💊 AI Drug Management (GLP-1 & Meds)": f"""Generate a comprehensive report on AI-powered diabetes drug
management for {p}. {d}
Sections: Executive Summary, GLP-1 Revolution (semaglutide, tirzepatide outcomes),
Smart Insulin Pens, AI Adherence Platforms, Once-Weekly Insulin, Digital Therapeutics
(Virta, Omada), Vendor Comparison, Recommendation, Safety & Limitations.""",

        "🥗 AI Nutrition & Lifestyle Coaching": f"""Generate a comprehensive report on AI nutrition and lifestyle
coaching for diabetes for {p}. {d}
Sections: Executive Summary, CGM-Guided Nutrition (DayTwo, Levels Health),
Meal Planning AI (One Drop, Nutrisense, Signos), Physical Activity AI,
Behavioral Change Platforms (Virta, Omada, Noom), App Comparison,
Clinical Evidence, Recommendation, Cost & Access.""",

        "📊 Population Health & Analytics": f"""Generate a comprehensive report on AI population health
management for diabetes for {p}. {d}
Sections: Executive Summary, Global Diabetes Burden, Risk Stratification AI
(Komodo, Clarify, Arcadia), CMS Prevention Programs, Payer Analytics,
Clinical Trial AI, Health Equity Considerations, ROI Evidence,
Vendor Comparison, Recommendation.""",

        "⚖️ Vendor Comparison Only": f"""Generate a detailed vendor comparison across ALL diabetes AI
categories for {p}. {d}
Create comparison tables for: AID Systems, CGM Devices, Diagnostic AI,
Drug Management AI, Nutrition Apps. Include: product, FDA status, outcomes,
price, target patient, strength, limitation.""",

        "🎯 Recommendation Only": f"""Provide direct, evidence-based AI recommendations for {p}. {d}
Structure: Quick Decision Framework, Recommendations by Sub-Category
(monitoring, delivery, diagnostics, drugs, nutrition), Complete AI Ecosystem
Recommendation (best combination of tools), Cost Considerations,
What to Avoid, Next Steps.""",
    }

    return module_prompts.get(module, module_prompts["🩺 Full Diabetes AI Platform"])

# ── JSON EXTRACTION PROMPT ────────────────────────────────────────
JSON_PROMPT = """Based on your report, return ONLY raw valid JSON (no markdown, no fences):
{
  "title": "AI in Diabetes Management: Intelligence Report",
  "subtitle": "Comprehensive Clinical Analysis & Recommendations",
  "module": "Full Diabetes AI Platform",
  "executive_summary": ["finding 1", "finding 2", "finding 3"],
  "key_metrics": [
    {"label": "Global Diabetes Burden", "value": "537M", "source": "IDF 2021"},
    {"label": "Projected by 2045", "value": "783M", "source": "IDF Diabetes Atlas"},
    {"label": "Annual Healthcare Cost", "value": "$966B", "source": "Global spend 2021"},
    {"label": "Best CGM Accuracy", "value": "7.9%", "source": "FreeStyle Libre 3 MARD"},
    {"label": "Best AID System TIR", "value": "78%", "source": "MiniMed 780G US users"},
    {"label": "T2D Reversal Rate", "value": "60%", "source": "Virta Health 1-year data"}
  ],
  "categories": [
    {"name": "AI Insulin Delivery", "top_product": "Omnipod 5", "key_outcome": "TIR 69%, hypo 1.12%", "fda": "Cleared T1D+T2D"},
    {"name": "AI Glucose Monitoring", "top_product": "FreeStyle Libre 3", "key_outcome": "MARD 7.9%, 14-day wear", "fda": "Cleared"},
    {"name": "AI Diagnostics", "top_product": "IDx-DR", "key_outcome": "87.2% sensitivity retinopathy", "fda": "Cleared"},
    {"name": "AI Drug Management", "top_product": "Tirzepatide + AI", "key_outcome": "HbA1c -2.3%, weight -22%", "fda": "Cleared"},
    {"name": "AI Nutrition", "top_product": "DayTwo", "key_outcome": "80% reduction glucose spikes", "fda": "N/A"},
    {"name": "Population Health AI", "top_product": "Komodo Health", "key_outcome": "18mo earlier identification", "fda": "N/A"}
  ],
  "recommendation": {
    "top_pick": "Omnipod 5 + FreeStyle Libre 3 + One Drop",
    "rationale": "Best combination for most diabetes patients: tubeless AID, most accurate affordable CGM, and proven digital coaching",
    "by_profile": [
      {"profile": "Type 1 — Active", "pick": "Omnipod 5 + Dexcom G7"},
      {"profile": "Type 2 — New to technology", "pick": "FreeStyle Libre 3 + Virta Health"},
      {"profile": "Prediabetes", "pick": "Omada Health + CGM screening"},
      {"profile": "Elderly", "pick": "FreeStyle Libre 3 + InPen"},
      {"profile": "Pediatric", "pick": "Omnipod 5 + CamAPS FX"}
    ]
  },
  "risks": [
    {"title": "Data Privacy", "desc": "Health data from multiple AI apps creates significant privacy exposure"},
    {"title": "Health Equity Gap", "desc": "AI tools disproportionately benefit insured, tech-savvy patients"},
    {"title": "Algorithm Bias", "desc": "HbA1c-based models underperform in Black patients by up to 30%"},
    {"title": "Cost & Access", "desc": "Full AI diabetes ecosystem costs $500-1200/month; coverage inconsistent"}
  ],
  "future_trends": [
    "Non-invasive glucose sensing (Apple Watch) mainstream by 2027-2028",
    "Fully closed-loop artificial pancreas without meal announcements by 2026",
    "AI-powered T2D reversal programs covered by Medicare by 2027",
    "Polypill + AI adherence reducing cardiovascular events by 40%",
    "Digital twin personalized diabetes management by 2029-2030"
  ]
}"""

# ── API CLIENT ────────────────────────────────────────────────────
def get_client():
    key = st.secrets.get("ANTHROPIC_API_KEY", os.environ.get("ANTHROPIC_API_KEY", ""))
    if not key:
        st.error("⚠️ No API key found. Add ANTHROPIC_API_KEY to Streamlit Secrets.")
        st.stop()
    return anthropic.Anthropic(api_key=key)

def generate_report(client, module, profile, depth, history):
    prompt   = build_prompt(module, profile, depth)
    system   = MODULE_SYSTEMS.get(module, MODULE_SYSTEMS["🩺 Full Diabetes AI Platform"])
    messages = history + [{"role": "user", "content": prompt}]
    response = client.messages.create(
        model="claude-opus-4-6",
        max_tokens=4000,
        system=system,
        messages=messages,
    )
    return response.content[0].text

def extract_slide_data(client, module, history):
    system   = MODULE_SYSTEMS.get(module, MODULE_SYSTEMS["🩺 Full Diabetes AI Platform"])
    messages = history + [{"role": "user", "content": JSON_PROMPT}]
    response = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=2500,
        system=system,
        messages=messages,
    )
    raw = response.content[0].text.strip()
    raw = re.sub(r'^```(?:json)?\s*', '', raw)
    raw = re.sub(r'\s*```$', '', raw)
    return json.loads(raw.strip())

# ── PPTX BUILDER ─────────────────────────────────────────────────
def build_pptx(d, module_name):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    # ── Slide 1: Title ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 10, 5.625, DARK)
    add_rect(s, 0, 0, 0.22, 5.625, ACCENT)
    add_rect(s, 0.22, 1.2, 9.78, 2.7, PRIMARY)
    add_text(s, "AI RESEARCH AGENT  ·  DIABETES TECHNOLOGY  ·  HCM 535",
             0.5, 0.75, 9, 0.32, 9, color=MID)
    add_text(s, d.get("title", "AI in Diabetes Management"),
             0.5, 1.32, 9.1, 1.1, 32, True, color=WHITE, font_name=FT)
    add_text(s, d.get("subtitle", ""), 0.5, 2.5, 9.1, 0.5,
             14, False, True, color=LIGHT)
    # Module badge
    mod_short = module_name.split(" ", 1)[1] if " " in module_name else module_name
    add_rect(s, 0.5, 3.15, 4.5, 0.38, MID)
    add_text(s, mod_short, 0.55, 3.15, 4.4, 0.38,
             11, True, color=DARK, align=PP_ALIGN.LEFT)
    # Student block
    add_rect(s, 0.5, 3.7, 5.6, 1.55, PRIMARY)
    info = [f"Student:     {STUDENT['name']}",
            f"Student ID:  {STUDENT['id']}",
            f"Institution: {STUDENT['school']}",
            f"Course:      {STUDENT['course']}"]
    for i, line in enumerate(info):
        add_text(s, line, 0.7, 3.8+i*0.34, 5.2, 0.32, 10, color=WHITE)
    add_rect(s, 6.5, 3.8, 2.8, 0.42, ACCENT)
    add_text(s, "Spring 2026", 6.5, 3.8, 2.8, 0.42,
             14, True, color=DARK, align=PP_ALIGN.CENTER, font_name=FT)

    # ── Slide 2: Executive Summary ──────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 10, 5.625, OFFWHITE)
    hdr(s, "Executive Summary", "AI in Diabetes — 2025 Evidence")
    add_rect(s, 0.3, 0.82, 4.5, 3.85, WHITE)
    add_text(s, "3 Key Findings", 0.5, 0.95, 4.1, 0.38,
             13, True, color=PRIMARY, font_name=FT)
    fcolors = [ACCENT, MID, PRIMARY]
    for i, txt in enumerate(d.get("executive_summary", ["Finding 1","Finding 2","Finding 3"])[:3]):
        y = 1.42 + i * 1.0
        add_rect(s, 0.5, y, 0.36, 0.36, fcolors[i])
        add_text(s, str(i+1), 0.5, y, 0.36, 0.36, 13, True,
                 color=DARK, align=PP_ALIGN.CENTER, font_name=FT)
        add_text(s, txt, 1.0, y, 3.6, 0.85, 10.5, color=TEXT)
    add_rect(s, 5.1, 0.82, 4.55, 1.65, DARK)
    add_text(s, "Global Diabetes Burden", 5.3, 0.9, 4.1, 0.32,
             11, True, color=MID, font_name=FT)
    add_text(s, "537 million people with diabetes globally (IDF 2021). "
                "Projected to reach 783 million by 2045. Annual healthcare "
                "cost: $966 billion. 50% remain undiagnosed worldwide.",
             5.3, 1.26, 4.1, 1.1, 10, color=LIGHT)
    add_rect(s, 5.1, 2.62, 4.55, 2.08, WHITE)
    add_text(s, "AI Impact on Diabetes Care", 5.3, 2.72, 4.0, 0.32,
             11, True, color=PRIMARY, font_name=FT)
    impacts = [("Time in Range improvement","+10–20%"),
               ("HbA1c reduction (AID)",    "−0.5–2.3%"),
               ("T2D reversal (Virta)",      "60% off insulin"),
               ("Earlier diagnosis",         "3–5 years sooner"),
               ("Cost savings/patient/yr",   "$2,650+")]
    for i,(lbl,val) in enumerate(impacts):
        y = 3.1+i*0.3
        add_text(s, lbl, 5.3, y, 2.5, 0.28, 9.5, color=GRAY)
        add_text(s, val, 7.8, y, 1.7, 0.28, 10, True, color=PRIMARY, align=PP_ALIGN.RIGHT)
    ftr(s, 2)

    # ── Slide 3: Key Metrics ────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 10, 5.625, OFFWHITE)
    hdr(s, "Key Metrics", "Clinical evidence 2024–2025")
    cw, ch, gx, gy = 3.0, 1.55, 0.12, 0.12
    sx, sy = 0.25, 0.92
    for i, m in enumerate(d.get("key_metrics", [])[:6]):
        col, row = i%3, i//3
        x = sx+col*(cw+gx); y = sy+row*(ch+gy)
        add_rect(s, x, y, cw, ch, WHITE)
        add_rect(s, x, y, cw, 0.05, MID)
        add_text(s, m.get("value","—"), x+0.15, y+0.1, cw-0.3, 0.72,
                 32, True, color=PRIMARY, font_name=FT)
        add_text(s, m.get("label",""), x+0.15, y+0.84, cw-0.3, 0.32,
                 11, True, color=TEXT)
        add_text(s, m.get("source",""), x+0.15, y+1.19, cw-0.3, 0.26,
                 8, False, True, color=GRAY)
    ftr(s, 3)

    # ── Slide 4: Categories Overview ────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 10, 5.625, DARK)
    hdr(s, "AI Categories in Diabetes Management", "Six pillars of AI-driven diabetes care")
    cat_colors = [MID, BLUE, PURPLE, rgb(0xB4,0x53,0x09), GREEN, PRIMARY]
    categories = d.get("categories", [])[:6]
    bw, bh = 2.95, 1.75
    positions = [(0.25,0.82),(3.38,0.82),(6.51,0.82),
                 (0.25,2.72),(3.38,2.72),(6.51,2.72)]
    for i,cat in enumerate(categories):
        if i >= len(positions): break
        x,y = positions[i]
        cc = cat_colors[i % len(cat_colors)]
        add_rect(s, x, y, bw, bh, PRIMARY)
        add_rect(s, x, y, bw, 0.05, cc)
        add_rect(s, x, y, 0.12, bh, cc)
        add_text(s, cat.get("name",""), x+0.2, y+0.1, bw-0.3, 0.38,
                 12, True, color=WHITE, font_name=FT)
        add_text(s, "Top: "+cat.get("top_product",""), x+0.2, y+0.55, bw-0.3, 0.3,
                 9.5, False, False, color=LIGHT)
        add_text(s, cat.get("key_outcome",""), x+0.2, y+0.88, bw-0.3, 0.5,
                 9, False, True, color=MID)
        add_rect(s, x+0.2, y+1.42, 1.5, 0.22, cc)
        add_text(s, cat.get("fda",""), x+0.2, y+1.42, 1.5, 0.22,
                 8, True, color=DARK, align=PP_ALIGN.CENTER)
    ftr(s, 4)

    # ── Slide 5: Recommendation ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 10, 5.625, OFFWHITE)
    hdr(s, "Recommendation", "Evidence-based selection guide")
    rec = d.get("recommendation", {})
    add_rect(s, 0.3, 0.82, 5.9, 2.15, DARK)
    add_rect(s, 0.3, 0.82, 0.18, 2.15, ACCENT)
    add_rect(s, 0.3, 0.82, 2.1, 0.32, ACCENT)
    add_text(s, "#1 TOP RECOMMENDATION", 0.32, 0.82, 2.05, 0.32,
             8, True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, rec.get("top_pick",""), 0.6, 1.18, 5.4, 0.6,
             16, True, color=WHITE, font_name=FT)
    add_text(s, rec.get("rationale",""), 0.6, 1.84, 5.3, 1.0,
             10, color=LIGHT)
    add_rect(s, 6.45, 0.82, 3.25, 4.12, WHITE)
    add_text(s, "By Patient Profile", 6.6, 0.9, 2.9, 0.34,
             12, True, color=PRIMARY, font_name=FT)
    pbg = [LIGHT, OFFWHITE, LIGHT, OFFWHITE, LIGHT]
    for i, p in enumerate(rec.get("by_profile",[])[:5]):
        y = 1.36+i*0.66
        add_rect(s, 6.55, y, 3.0, 0.58, pbg[i%2])
        add_text(s, p.get("profile",""), 6.68, y+0.04, 2.8, 0.22,
                 8, False, True, color=GRAY)
        add_text(s, p.get("pick",""), 6.68, y+0.27, 2.8, 0.27,
                 10.5, True, color=PRIMARY, font_name=FT)
    ftr(s, 5)

    # ── Slide 6: Risks ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 10, 5.625, OFFWHITE)
    hdr(s, "Risks & Limitations", "Critical considerations")
    rcolors = [RED, rgb(0x85,0x4D,0x0E), PURPLE, PRIMARY]
    for i,r in enumerate(d.get("risks",[])[:4]):
        col,row = i%2, i//2
        x=0.3+col*4.9; y=0.92+row*2.12
        add_rect(s, x, y, 4.55, 1.97, WHITE)
        add_rect(s, x, y, 4.55, 0.05, rcolors[i])
        add_rect(s, x, y, 0.18, 1.97, rcolors[i])
        add_text(s, r.get("title",""), x+0.28, y+0.12, 4.1, 0.38,
                 13, True, color=TEXT, font_name=FT)
        add_text(s, r.get("desc",""), x+0.28, y+0.55, 4.1, 1.28, 11, color=GRAY)
    ftr(s, 6)

    # ── Slide 7: Future Outlook ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 10, 5.625, DARK)
    hdr(s, "Future Outlook 2025–2030", "Next generation AI diabetes care")
    years=["2026","2027","2027","2028","2029–30"]
    bw3,bh3=1.72,3.05; sx3,sy3,gap3=0.22,1.0,0.1
    for i,t in enumerate(d.get("future_trends",[])[:5]):
        x=sx3+i*(bw3+gap3)
        add_rect(s, x, sy3, bw3, bh3, PRIMARY)
        add_rect(s, x, sy3, bw3, 0.05, ACCENT)
        add_rect(s, x+0.1, sy3+0.12, 1.5, 0.3, ACCENT)
        add_text(s, years[i], x+0.1, sy3+0.12, 1.5, 0.3,
                 10, True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(s, t, x+0.1, sy3+0.55, bw3-0.2, 2.35, 10, color=WHITE)
    ftr(s, 7)

    # ── Slide 8: Closing ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 10, 5.625, DARK)
    add_rect(s, 0, 0, 0.22, 5.625, ACCENT)
    add_text(s, "Thank You", 0.5, 0.55, 9, 0.9, 40, True, color=WHITE, font_name=FT)
    add_text(s, "AI in Diabetes Management: Clinical Evidence & Strategic Recommendation",
             0.5, 1.52, 9, 0.45, 13, False, True, color=LIGHT)
    add_rect(s, 0.5, 2.15, 5.5, 1.6, PRIMARY)
    rows=[("Student",STUDENT["name"]),("Student ID",STUDENT["id"]),
          ("Institution",STUDENT["school"]),("Course",STUDENT["course"])]
    for i,(lbl,val) in enumerate(rows):
        add_text(s, lbl+":", 0.7, 2.27+i*0.35, 1.4, 0.3, 10, True, color=MID)
        add_text(s, val, 2.15, 2.27+i*0.35, 3.7, 0.3, 10, color=WHITE)
    add_rect(s, 0.5, 3.92, 9.0, 1.4, PRIMARY)
    add_text(s, "Key References", 0.7, 3.98, 3.0, 0.28, 10, True, color=MID, font_name=FT)
    refs=("IDF Diabetes Atlas 10th Edition (2021). Global diabetes statistics.\n"
          "Forlenza et al. (2024). Real-World Evidence of Omnipod 5, n=69,902. Diabetes Technol Ther.\n"
          "Mohanadas et al. (2026). AI in Medical Devices. JMIR 28:e72410.\n"
          "Virta Health (2024). T2D Reversal Program 5-Year Outcomes.")
    add_text(s, refs, 0.7, 4.28, 8.6, 1.0, 8, color=LIGHT)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

# ── SESSION STATE ─────────────────────────────────────────────────
for k,v in [("history",[]),("report_text",""),("slide_data",None),
            ("pptx_buffer",None),("pptx_ready",False),
            ("authenticated",False),("current_module","")]:
    if k not in st.session_state:
        st.session_state[k] = v

# ── ACCESS GATE ──────────────────────────────────────────────────
if not st.session_state.authenticated:
    st.markdown("""
    <div class="main-header">
        <h1>🩺 Diabetes AI Agent</h1>
        <p>AI Applications in Diabetes Management — Full Platform</p>
        <div class="student-badge">Hugo Silva  ·  ID 74964557  ·  ITU SPRING 2026  ·  HCM 535</div>
    </div>""", unsafe_allow_html=True)
    c1,c2,c3 = st.columns([1,2,1])
    with c2:
        st.markdown("### Enter Access Code")
        code = st.text_input("Access code", type="password",
                             placeholder="Enter code to continue...")
        if st.button("Unlock Agent"):
            if code == "HCM535":
                st.session_state.authenticated = True
                st.rerun()
            else:
                st.error("Incorrect code.")
    st.stop()

# ── SIDEBAR ───────────────────────────────────────────────────────
with st.sidebar:
    st.markdown("### 🔬 Research Module")
    module = st.selectbox("Select module", list(MODULES.keys()))
    info = MODULES[module]
    st.markdown(f"""
    <div class="module-card">
        <h4>{info['icon']} {info['tag']}</h4>
        <p>{info['desc']}</p>
    </div>""", unsafe_allow_html=True)

    st.markdown("### ⚙️ Settings")
    profile = st.selectbox("Patient profile", [
        "General","Type 1 Diabetes","Type 2 Diabetes",
        "Prediabetes / At-Risk","Pediatric","Elderly (65+)",
        "Newly Diagnosed","Healthcare Provider"
    ])
    depth = st.selectbox("Report depth", [
        "Executive Summary","Detailed Analysis","Clinical Deep-Dive"
    ])

    st.markdown("---")
    st.markdown("### 💬 Follow-up")
    followup = st.text_area("Ask the agent", placeholder="e.g. Which CGM is best for a child?", height=80)
    ask_btn = st.button("Ask ↗")

    st.markdown("---")
    st.markdown(f"""
    <div style="background:#F1F5F9;border-radius:10px;padding:0.8rem;font-size:0.78rem;color:#334155">
    <strong>Modules available:</strong><br>
    {'<br>'.join([f"{v['icon']} {k.split(' ',1)[1] if ' ' in k else k}" for k,v in MODULES.items()])}
    </div>""", unsafe_allow_html=True)

    if st.button("🔄 Reset"):
        for k in ["history","report_text","slide_data","pptx_buffer","pptx_ready"]:
            st.session_state[k] = [] if k=="history" else None if k in ["slide_data","pptx_buffer"] else "" if k=="report_text" else False
        st.rerun()

# ── MAIN ──────────────────────────────────────────────────────────
st.markdown("""
<div class="main-header">
    <h1>🩺 Diabetes AI Agent</h1>
    <p>AI Applications in Diabetes Management — 6 Research Modules · Clinical Intelligence · PowerPoint Generator</p>
    <div class="student-badge">Hugo Silva  ·  ID 74964557  ·  ITU SPRING 2026  ·  HCM 535 – Data Analytics in Healthcare</div>
</div>""", unsafe_allow_html=True)

# Metrics row
c1,c2,c3,c4 = st.columns(4)
metrics = [
    ("537M","People with diabetes globally"),
    ("6","AI research modules"),
    ("78%","Best AID system TIR"),
    ("60%","T2D reversal rate (Virta)"),
]
for col,(val,lbl) in zip([c1,c2,c3,c4],metrics):
    with col:
        st.markdown(f'<div class="metric-card"><h3>{val}</h3><p>{lbl}</p></div>',
                    unsafe_allow_html=True)

st.markdown('<hr class="teal-divider">', unsafe_allow_html=True)

# Module display
st.markdown(f"**Selected module:** {module} — *{MODULES[module]['desc']}*")

# Generate button
col_btn, col_status = st.columns([2,3])
with col_btn:
    gen_btn = st.button("▶  Generate Report + PPT", use_container_width=True)
with col_status:
    if st.session_state.report_text:
        st.markdown('<span class="badge-success">✓ Report ready</span>',
                    unsafe_allow_html=True)

# ── GENERATION PIPELINE ───────────────────────────────────────────
if gen_btn:
    client  = get_client()
    prog    = st.progress(0)
    status  = st.empty()
    try:
        status.markdown(f"**Step 1/3** — Generating {module} report...")
        prog.progress(10)
        report = generate_report(client, module, profile, depth,
                                 st.session_state.history)
        st.session_state.report_text    = report
        st.session_state.current_module = module
        st.session_state.history.append({"role":"user",    "content":f"Generate: {module}, {profile}, {depth}"})
        st.session_state.history.append({"role":"assistant","content":report})
        prog.progress(45)

        status.markdown("**Step 2/3** — Extracting structured slide data...")
        slide_data = extract_slide_data(client, module, st.session_state.history)
        st.session_state.slide_data = slide_data
        prog.progress(75)

        status.markdown("**Step 3/3** — Building PowerPoint presentation...")
        buf = build_pptx(slide_data, module)
        st.session_state.pptx_buffer = buf
        st.session_state.pptx_ready  = True
        prog.progress(100)
        status.markdown("✅ **Done!** Report and PPT ready.")

    except Exception as e:
        st.error(f"Error: {e}")
        prog.empty(); status.empty()

# ── FOLLOW-UP ─────────────────────────────────────────────────────
if ask_btn and followup.strip():
    if not st.session_state.history:
        st.warning("Generate a report first.")
    else:
        client = get_client()
        with st.spinner("Thinking..."):
            try:
                system = MODULE_SYSTEMS.get(
                    st.session_state.current_module,
                    MODULE_SYSTEMS["🩺 Full Diabetes AI Platform"]
                )
                msgs = st.session_state.history + [{"role":"user","content":followup}]
                resp = client.messages.create(
                    model="claude-opus-4-6", max_tokens=1500,
                    system=system, messages=msgs
                )
                ans = resp.content[0].text
                st.session_state.history.append({"role":"user",    "content":followup})
                st.session_state.history.append({"role":"assistant","content":ans})
                st.session_state.report_text += f"\n\n---\n\n**Q: {followup}**\n\n{ans}"
            except Exception as e:
                st.error(f"Follow-up error: {e}")

# ── RESULTS ───────────────────────────────────────────────────────
if st.session_state.report_text:
    tab1, tab2 = st.tabs(["📄 Research Report", "📊 Slide Data"])
    with tab1:
        st.markdown(st.session_state.report_text)
    with tab2:
        if st.session_state.slide_data:
            st.json(st.session_state.slide_data)

    st.markdown('<hr class="teal-divider">', unsafe_allow_html=True)
    dl1, dl2 = st.columns(2)
    ts = datetime.now().strftime('%Y%m%d_%H%M')
    with dl1:
        st.download_button(
            label="⬇️  Download Report (.txt)",
            data=st.session_state.report_text.encode(),
            file_name=f"diabetes_AI_report_{ts}.txt",
            mime="text/plain",
            use_container_width=True,
        )
    with dl2:
        if st.session_state.pptx_ready and st.session_state.pptx_buffer:
            st.download_button(
                label="⬇️  Download Presentation (.pptx)",
                data=st.session_state.pptx_buffer,
                file_name=f"diabetes_AI_agent_{ts}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )

# ── FOOTER ────────────────────────────────────────────────────────
st.markdown(f"""
<div class="footer">
    Hugo Silva  ·  ID 74964557  ·  ITU | SPRING 2026  ·  HCM 535 – Data Analytics Application in Healthcare<br>
    Powered by Claude (Anthropic)  ·  Not a substitute for medical advice  ·  Sources: IDF 2021, JMIR 2026, FDA Device Database
</div>""", unsafe_allow_html=True)
